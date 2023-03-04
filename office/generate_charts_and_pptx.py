import os
import pandas as pd
import seaborn as sns
from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches


def generate_charts_and_pptx(folder_path):
    try:
        # 检查"charts"文件夹是否存在，如果不存在则创建
        if not os.path.exists('charts'):
            os.makedirs('charts')

        # 获取文件夹中的所有Excel文件
        excel_files = [f for f in os.listdir(folder_path) if f.endswith('.xlsx')]

        # 创建PPTX文档
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]

        # 遍历每个Excel文件并创建图表和幻灯片
        for excel_file in excel_files:
            # 读取Excel文件中的数据
            file_path = os.path.join(folder_path, excel_file)
            df = pd.read_excel(file_path, sheet_name=0)

            # 按产品分组并计算销售总额
            sales_by_product = df.dropna().groupby('Product')['Sales'].sum().reset_index()

            # 创建图表并保存为文件
            chart_title = os.path.splitext(excel_file)[0]
            chart_filename = os.path.join('charts', '{}.png'.format(chart_title))

            sns.set(style="whitegrid")
            ax = sns.barplot(x="Product", y="Sales", data=sales_by_product)
            ax.set_title(chart_title)
            ax.set_xlabel('Product')
            ax.set_ylabel('Sales')

            fig = ax.get_figure()
            fig.savefig(chart_filename)
            plt.close(fig)

            # 在PPTX中创建新的幻灯片
            slide = prs.slides.add_slide(title_slide_layout)
            slide_title = slide.shapes.title
            slide_title.text = chart_title

            # 将图表插入幻灯片
            left = Inches(1)
            top = Inches(2)
            height = Inches(6)
            pic = slide.shapes.add_picture(chart_filename, left, top, height=height)

        # 将幻灯片保存为PPTX文件
        pptx_filename = os.path.join(folder_path, 'financial_data.pptx')
        prs.save(pptx_filename)

        print('Charts and PowerPoint file generated successfully.')

    except Exception as e:
        print('Error: {}'.format(e))


folder_path = '/Users/alexwang/Projects/pythonUtil/office/data'

if os.path.exists(folder_path):
    generate_charts_and_pptx(folder_path)
else:
    print('Error: The specified folder does not exist.')
